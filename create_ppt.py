#!/usr/bin/env python3
"""
自动创建PowerPoint演讲文稿
运行: python3 create_ppt.py
输出: NTD_Presentation.pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("正在安装 python-pptx...")
    import subprocess
    subprocess.run(['pip', 'install', 'python-pptx'], check=True)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

import os

print("创建PowerPoint演讲文稿...")

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 比例
prs.slide_height = Inches(5.625)

# 定义布局
blank_layout = prs.slide_layouts[6]  # 空白布局
title_layout = prs.slide_layouts[0]  # 标题布局
title_content = prs.slide_layouts[1]  # 标题+内容

# 颜色定义
BLUE = RGBColor(52, 152, 219)
DARK = RGBColor(44, 62, 80)
RED = RGBColor(231, 76, 60)
GREEN = RGBColor(39, 174, 96)

def add_title_slide(title, subtitle):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(title, content_list):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(title_content)
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_image_slide(title, image_path, caption=""):
    """添加带图片的幻灯片"""
    slide = prs.slides.add_slide(blank_layout)

    # 添加标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE

    # 添加图片
    if os.path.exists(image_path):
        left = Inches(1.5)
        top = Inches(1.2)
        height = Inches(3.5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        # 如果图片不存在，添加提示
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
        tf = txBox.text_frame
        tf.text = f"[图片: {os.path.basename(image_path)}]"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if caption:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.4))
        tf = txBox.text_frame
        tf.text = caption
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)

    return slide

print("添加幻灯片...")

# Slide 1: 标题
add_title_slide(
    "Investigating Trends in Neglected Tropical Disease (NTD) Clinical Studies",
    "Group 16 - Lancaster University\nWHO ICTRP Data 1999-2023"
)

# Slide 2: 项目概览
add_content_slide("Project Overview", [
    "Source: WHO ICTRP",
    "Period: 1999-2023 (25 years)",
    "Trials: 311 valid trials",
    "Countries: 62",
    "Diseases: Chagas, Schistosomiasis, Endemic worms, Visceral leishmaniasis"
])

# Slide 3: 研究问题
add_content_slide("Research Questions", [
    "RQ 2.2.1: What factors affect results publication?",
    "RQ 2.2.2: Can network analysis reveal country partnerships?",
    "RQ 2.2.3: What proportion funded by pharma companies?",
    "RQ 2.2.4: Are children and pregnant women included?",
    "RQ 2.2.5: Which drugs studied for Chagas disease?"
])

# Slide 4: RQ 2.2.1 - 发表因素
add_content_slide("RQ 2.2.1: Publication Factors", [
    "FINDING: Severe Publication Gap",
    "",
    "Publication rate: 4.2% (13/311)",
    "Unpublished: 95.8% (298/311)",
    "",
    "Method: Logistic Regression",
    "Features: phase, study type, sponsor, income level",
    "",
    "Implication: Severe publication bias"
])

# Slide 5: 发表预测因子（带图）
add_image_slide(
    "RQ 2.2.1: Publication Predictors",
    "CleanedDataPlt/coefficients_plot.png",
    "Factors increasing (+) and decreasing (-) publication likelihood"
)

# Slide 6: 制药资金（带图）
add_image_slide(
    "RQ 2.2.3: Pharmaceutical Funding",
    "CleanedDataPlt/sponsor_distribution.png",
    "Non-profit: 66.9% | Industry: 7.4% | Government: 4.5%"
)

# Slide 7: 工业与疾病负担（带图）
add_image_slide(
    "RQ 2.2.3: Industry vs Disease Burden",
    "CleanedDataPlt/industry_region.png",
    "60% of industry trials NOT in high-burden regions"
)

# Slide 8: 国家合作
add_content_slide("RQ 2.2.2: Country Collaborations", [
    "FINDING: Strong International Partnerships",
    "",
    "Multi-country trials: 43",
    "Countries in network: 43",
    "Collaborative connections: 163",
    "",
    "Top Hub Countries:",
    "1. Argentina (22 partners, Centrality: 0.294)",
    "2. Cambodia (18 partners, Centrality: 0.275)",
    "3. Brazil (18 partners, Centrality: 0.159)"
])

# Slide 9: 特殊人群
add_content_slide("RQ 2.2.4: Special Populations", [
    "FINDING: Limited Inclusion",
    "",
    "Children:",
    "  • No explicit tracking field",
    "  • Estimated ~50% may include children",
    "",
    "Pregnant Women:",
    "  • Explicitly included: 12 trials (3.8%)",
    "  • Excluded/unclear: 96.2%",
    "",
    "Implication: Vulnerable populations underrepresented"
])

# Slide 10: Chagas药物
add_content_slide("RQ 2.2.5: Drug Trends (Chagas)", [
    "FINDING: Benznidazole Dominates",
    "",
    "Total Chagas trials: 88",
    "Method: Regex extraction",
    "",
    "Top 5 Drugs:",
    "1. Benznidazole - 21 trials",
    "2. Placebo - 10 trials",
    "3. E1224 - 6 trials",
    "4. Nifurtimox - 5 trials",
    "",
    "Implication: Limited drug diversity"
])

# Slide 11: 如何解决研究目标
add_content_slide("How Results Address Research Goals", [
    "Goal: Clarify evolution of NTD trial activities",
    "",
    "✓ Temporal Trends: 25-year coverage",
    "✓ Geographic Distribution: 62 countries analyzed",
    "✓ Funding Structure: Non-profit 67%, Pharma 7%",
    "✓ Publication Patterns: 4.2% gap identified",
    "✓ Collaboration: Network analysis reveals hubs"
])

# Slide 12: 关键发现
add_content_slide("Key Discoveries", [
    "3 Critical Findings:",
    "",
    "1. PUBLICATION CRISIS",
    "   → 95.8% unpublished",
    "",
    "2. FUNDING GAP",
    "   → Pharma 7.4%, Non-profits 67%",
    "",
    "3. REPRESENTATION GAP",
    "   → Pregnant women 3.8%, Children unclear"
])

# Slide 13: 数据局限性
add_content_slide("Data Limitations", [
    "Challenges:",
    "",
    "• Severe publication bias (only 13 published)",
    "• Data completeness (4 outliers, missing values)",
    "• Classification accuracy (keyword-based)",
    "• Temporal imbalance (varies by year)"
])

# Slide 14: 方法局限性
add_content_slide("Methodological Limitations", [
    "Analytical Challenges:",
    "",
    "1. Logistic Regression:",
    "   • Small sample (n=13)",
    "   • Class imbalance (4% vs 96%)",
    "",
    "2. Network Analysis:",
    "   • Country-level only",
    "",
    "3. Drug Extraction:",
    "   • Regex-based, may miss variations",
    "",
    "Mitigation: Stratified sampling, transparent docs"
])

# Slide 15: 改进建议
add_content_slide("Recommendations", [
    "Data Collection:",
    "• Standardize special population fields",
    "• Mandate result publication",
    "",
    "Methodological:",
    "• Handle class imbalance (SMOTE)",
    "• Try ensemble methods",
    "",
    "Policy:",
    "• Mandate publication (address 96% gap)",
    "• Incentivize pharma investment",
    "• Target high-burden regions"
])

# Slide 16: 分析优势
add_content_slide("Strengths of Our Analysis", [
    "✓ Comprehensive Cleaning",
    "  Robust outlier detection, automated classification",
    "",
    "✓ Multi-Method Approach",
    "  Statistics + ML + Network + Text Mining",
    "",
    "✓ Reproducible Pipeline",
    "  5 modular scripts, complete documentation",
    "",
    "✓ Transparent & Practical",
    "  All limitations acknowledged, actionable insights"
])

# Slide 17: 要点总结
add_content_slide("Key Takeaways", [
    "NTD research needs structural reform:",
    "• Better publication incentives",
    "• Increased pharma engagement",
    "• Improved data standards",
    "• Inclusive trial design",
    "",
    "Our Contribution:",
    "• Evidence for policy makers",
    "• Guidance for funders",
    "• Roadmap for researchers",
    "",
    "311 trials → 5 questions → 3 critical gaps"
])

# Slide 18: 结束
add_title_slide(
    "Thank You",
    "Questions?\n\nGroup 16 - Lancaster University\nPeizhe Jiang • Congyao Ren • Zixu Wang • Alaghwani Balsam"
)

# 保存文件
output_file = "NTD_Presentation.pptx"
prs.save(output_file)

print(f"\n✓ PowerPoint已创建: {output_file}")
print(f"✓ 共 {len(prs.slides)} 张幻灯片")
print("\n可以用Microsoft PowerPoint或Google Slides打开！")
